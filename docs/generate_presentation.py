r"""
Generate a professional PowerPoint deck for the Text-to-SQL project.

Themed to match the Giesecke+Devrient "EOC Townhall" corporate template:
  - Deep navy background (#0E2841)
  - Azure-blue primary accent (#0F9ED5) with a brighter cyan and yellow highlight
  - White "Giesecke+Devrient / Creating Confidence" logo footer on every slide

Run with the project venv (python-pptx installed):
    & .\.env\Scripts\python.exe .\docs\generate_presentation.py
or the base Python:
    & "C:\\Users\\gangwarr\\AppData\\Local\\Programs\\Python\\Python314\\python.exe" docs/generate_presentation.py

Output: docs/TextToSQL_Presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ----------------------------------------------------------------------------
# Theme  —  Giesecke+Devrient corporate palette (from EOC Townhall reference)
# ----------------------------------------------------------------------------
BG_DARK = RGBColor(0x0E, 0x28, 0x41)    # G+D deep navy background
BG_CARD = RGBColor(0x16, 0x39, 0x5A)    # card surface (lighter navy)
BG_CARD2 = RGBColor(0x1F, 0x4D, 0x73)   # lighter card surface
TEAL = RGBColor(0x0F, 0x9E, 0xD5)       # PRIMARY accent = G+D azure blue
TEAL_DK = RGBColor(0x0B, 0x7C, 0xAA)    # darker azure
CYAN = RGBColor(0x35, 0xC5, 0xE8)       # bright cyan (title gradient)
WHITE = RGBColor(0xF2, 0xF7, 0xFC)      # near-white text
GRAY = RGBColor(0xAF, 0xC2, 0xD4)       # muted light-blue-gray text
BLUE = RGBColor(0x2B, 0xB8, 0xE6)       # bright cyan-blue accent
PURPLE = RGBColor(0xB5, 0x6A, 0xCE)     # G+D purple (lightened for navy)
AMBER = RGBColor(0xF0, 0x88, 0x3E)      # G+D orange accent
GREEN = RGBColor(0x57, 0xB8, 0x47)      # G+D green accent
RED = RGBColor(0xE8, 0x51, 0x55)        # danger / problem
PINK = RGBColor(0xD8, 0x6B, 0xB0)       # secondary accent
YELLOW = RGBColor(0xFF, 0xD2, 0x3F)     # highlight (G+D active-item yellow)

FONT = "Segoe UI"                       # closest installed match to Aptos
FONT_LIGHT = "Segoe UI Light"
MONO = "Consolas"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

# G+D logo (white "Giesecke+Devrient / Creating Confidence")
ASSETS_DIR = Path(__file__).resolve().parent / "assets"
LOGO_PATH = ASSETS_DIR / "gd_logo_white.png"


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def set_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    no_line(shape)
    shape.shadow.inherit = False


def add_rect(slide, left, top, width, height, color, rounded=False, radius=0.08):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    fill_shape(shp, color)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def add_outline_rect(slide, left, top, width, height, line_color, rounded=True, radius=0.06, weight=1.25):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.background()
    shp.line.color.rgb = line_color
    shp.line.width = Pt(weight)
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, wrap=True, space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (text, size, color, bold, font) tuples."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold, font) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def one(text, size, color, bold=False, font=FONT):
    """Shorthand to build a single-run paragraph."""
    return [(text, size, color, bold, font)]


def fit_text_in_shape(shape, runs, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT,
                      ml=0.12, mr=0.12, mt=0.06, mb=0.06, line_spacing=1.0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold, font) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font


def add_footer(slide, number=None, logo_height=0.32):
    """Place the white G+D logo bottom-left and an optional page number bottom-right."""
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.55), Inches(7.0), height=Inches(logo_height))
    else:
        # Fallback to text if the logo asset is missing
        add_text(slide, Inches(0.55), Inches(6.98), Inches(5.0), Inches(0.4),
                 [[("Giesecke+Devrient", 11, WHITE, True, FONT),
                   ("   Creating Confidence", 9.5, GRAY, False, FONT)]])
    if number is not None:
        add_text(slide, Inches(12.1), Inches(7.02), Inches(0.95), Inches(0.35),
                 [one(f"{number:02d}", 11, GRAY, bold=True)], align=PP_ALIGN.RIGHT)


def header(slide, kicker, title, number=None):
    """Standard slide header: accent bar + kicker + title + G+D footer/logo."""
    # accent vertical bar
    add_rect(slide, Inches(0.55), Inches(0.55), Inches(0.11), Inches(0.95), TEAL, rounded=True, radius=0.5)
    add_text(slide, Inches(0.82), Inches(0.5), Inches(11.0), Inches(0.4),
             [one(kicker.upper(), 12.5, TEAL, bold=True)])
    add_text(slide, Inches(0.8), Inches(0.82), Inches(11.6), Inches(0.8),
             [one(title, 30, WHITE, bold=True)])
    add_footer(slide, number)


def chip(slide, left, top, text, color, width=None):
    w = width or Inches(0.46 + 0.108 * len(text))
    box = add_rect(slide, left, top, w, Inches(0.36), BG_CARD2, rounded=True, radius=0.5)
    box.line.color.rgb = color
    box.line.width = Pt(1)
    fit_text_in_shape(box, [one(text, 11, color, bold=True)], anchor=MSO_ANCHOR.MIDDLE,
                      align=PP_ALIGN.CENTER, ml=0.08, mr=0.08)
    return w


def bullet_block(slide, left, top, width, items, gap=0.62, dot=TEAL, title_size=15, desc_size=11.5):
    """items: list of (title, desc) or (title, desc, color)."""
    y = top
    for it in items:
        title = it[0]
        desc = it[1] if len(it) > 1 else ""
        c = it[2] if len(it) > 2 else dot
        add_rect(slide, left, y + Inches(0.06), Inches(0.16), Inches(0.16), c, rounded=True, radius=0.5)
        paras = [[(title, title_size, WHITE, True, FONT)]]
        if desc:
            paras.append([(desc, desc_size, GRAY, False, FONT)])
        add_text(slide, left + Inches(0.34), y, width - Inches(0.34), Inches(gap),
                 paras, space_after=2, line_spacing=1.0)
        y += Inches(gap)


def hero_logo(slide):
    """Large G+D logo top-right for title / closing hero slides."""
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(9.85), Inches(0.5), height=Inches(0.46))


# ----------------------------------------------------------------------------
# Slides
# ----------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    def new_slide(color=BG_DARK):
        s = prs.slides.add_slide(blank)
        set_bg(s, color)
        return s

    # ---- Slide 1: Title ----------------------------------------------------
    s = new_slide()
    # decorative cyan/navy panels (G+D gradient feel)
    add_rect(s, Inches(0), Inches(0), Inches(0.28), EMU_H, TEAL)
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.4), Inches(-1.4), Inches(5.2), Inches(5.2))
    fill_shape(c1, BG_CARD)
    c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(3.6), Inches(4.2), Inches(4.2))
    fill_shape(c2, BG_CARD2)
    hero_logo(s)

    chip(s, Inches(0.9), Inches(1.45), "FULL-STACK  •  AI  •  POSTGRESQL", TEAL, width=Inches(3.7))
    add_text(s, Inches(0.85), Inches(2.1), Inches(11.5), Inches(2.0),
             [[("Text", 66, WHITE, True, FONT), ("-to-", 66, GRAY, True, FONT), ("SQL", 66, TEAL, True, FONT)]])
    add_text(s, Inches(0.9), Inches(3.6), Inches(10.8), Inches(1.0),
             [one("Turn plain-English questions into accurate SQL — and answers.", 22, WHITE, bold=False)])
    add_text(s, Inches(0.9), Inches(4.4), Inches(10.6), Inches(1.0),
             [one("A full-stack application that uses Azure OpenAI, semantic search and "
                  "join-graph reasoning to generate, validate and run SQL on PostgreSQL.", 14.5, GRAY)],
             line_spacing=1.2)

    # tech chips row
    x = Inches(0.9)
    for label, col in [("React 18", BLUE), ("FastAPI", GREEN), ("Azure OpenAI GPT-4.1", PURPLE),
                       ("FAISS", AMBER), ("PostgreSQL", TEAL)]:
        w = chip(s, x, Inches(5.55), label, col)
        x += w + Inches(0.18)

    add_footer(s)

    # ---- Slide 2: Overview -------------------------------------------------
    s = new_slide()
    header(s, "Introduction", "Project Overview", 2)
    add_text(s, Inches(0.85), Inches(1.85), Inches(7.1), Inches(1.4),
             [one("Text-to-SQL lets non-technical users query a relational database by simply "
                  "asking a question in natural language. The system finds the right tables, "
                  "writes the SQL, runs it safely, and explains the result.", 15, GRAY)],
             line_spacing=1.3)

    items = [
        ("Conversational querying", "No SQL knowledge required — just ask.", BLUE),
        ("Retrieval-augmented (RAG)", "Embeddings pick only the relevant tables.", PURPLE),
        ("Schema & join aware", "Uses foreign-key relationships to join correctly.", TEAL),
        ("Safe by design", "Detects and blocks destructive operations.", RED),
    ]
    bullet_block(s, Inches(0.85), Inches(3.35), Inches(7.2), items, gap=0.78)

    # right stat cards
    cards = [("3", "Application tiers", BLUE), ("5", "AI pipeline steps", TEAL),
             ("2", "LLM calls / query", PURPLE), ("100%", "Read-only execution", GREEN)]
    cx, cy = Inches(8.35), Inches(1.9)
    for i, (big, label, col) in enumerate(cards):
        col_i = i % 2
        row_i = i // 2
        left = cx + col_i * Inches(2.25)
        top = cy + row_i * Inches(2.25)
        card = add_rect(s, left, top, Inches(2.05), Inches(2.05), BG_CARD, rounded=True, radius=0.1)
        add_rect(s, left, top, Inches(2.05), Inches(0.12), col, rounded=False)
        add_text(s, left, top + Inches(0.45), Inches(2.05), Inches(0.9),
                 [one(big, 40, col, bold=True)], align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.1), top + Inches(1.42), Inches(1.85), Inches(0.5),
                 [one(label, 11.5, GRAY, bold=True)], align=PP_ALIGN.CENTER)

    # ---- Slide 3: Problem & Solution --------------------------------------
    s = new_slide()
    header(s, "Motivation", "The Problem & Our Solution", 3)
    # Problem card
    pcard = add_rect(s, Inches(0.85), Inches(1.95), Inches(5.6), Inches(4.4), BG_CARD, rounded=True, radius=0.06)
    add_rect(s, Inches(0.85), Inches(1.95), Inches(5.6), Inches(0.6), RED, rounded=True, radius=0.2)
    add_text(s, Inches(1.05), Inches(2.0), Inches(5.2), Inches(0.5), [one("THE PROBLEM", 14, WHITE, bold=True)],
             anchor=MSO_ANCHOR.MIDDLE)
    prob = [
        ("SQL is a barrier", "Business users can't write queries themselves.", RED),
        ("Schemas are complex", "Dozens of tables, keys and join paths to remember.", RED),
        ("Slow feedback loop", "Every question needs an engineer in the loop.", RED),
        ("Risky operations", "A wrong query can modify or delete data.", RED),
    ]
    bullet_block(s, Inches(1.1), Inches(2.85), Inches(5.1), prob, gap=0.82, title_size=13.5, desc_size=11)

    # Solution card
    scard = add_rect(s, Inches(6.9), Inches(1.95), Inches(5.55), Inches(4.4), BG_CARD, rounded=True, radius=0.06)
    add_rect(s, Inches(6.9), Inches(1.95), Inches(5.55), Inches(0.6), GREEN, rounded=True, radius=0.2)
    add_text(s, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.5), [one("OUR SOLUTION", 14, BG_DARK, bold=True)],
             anchor=MSO_ANCHOR.MIDDLE)
    sol = [
        ("Ask in plain English", "Natural-language input, instant SQL + results.", GREEN),
        ("AI finds the tables", "Semantic search retrieves only what's relevant.", GREEN),
        ("Automatic, correct joins", "Join graph guides multi-table queries.", GREEN),
        ("Guard-railed execution", "Dangerous operations detected and blocked.", GREEN),
    ]
    bullet_block(s, Inches(7.15), Inches(2.85), Inches(5.05), sol, gap=0.82, title_size=13.5, desc_size=11)

    # ---- Slide 4: Key Features --------------------------------------------
    s = new_slide()
    header(s, "Capabilities", "Key Features", 4)
    feats = [
        ("Natural Language → SQL", "Azure OpenAI GPT-4.1 writes valid PostgreSQL from a question.", BLUE),
        ("Semantic Table Retrieval", "Embeddings + FAISS pick the top-k relevant tables.", PURPLE),
        ("Join-Graph Reasoning", "Foreign-key graph builds correct multi-table joins.", TEAL),
        ("Danger Detection", "Blocks DELETE / DROP / UPDATE / INSERT before execution.", RED),
        ("CSV / Excel Upload", "Drop a file → instantly query it as a table.", AMBER),
        ("AI Result Summaries", "A second LLM call explains the results in plain words.", GREEN),
        ("Query Complexity Score", "Rates each query Simple → Very Complex.", PINK),
        ("Multi-Database Support", "Browse schemas and switch databases on the fly.", CYAN),
    ]
    gx, gy = Inches(0.85), Inches(1.95)
    cw, ch = Inches(5.75), Inches(1.12)
    for i, (title, desc, col) in enumerate(feats):
        col_i = i % 2
        row_i = i // 2
        left = gx + col_i * (cw + Inches(0.3))
        top = gy + row_i * (ch + Inches(0.12))
        card = add_rect(s, left, top, cw, ch, BG_CARD, rounded=True, radius=0.1)
        add_rect(s, left, top, Inches(0.12), ch, col, rounded=False)
        add_text(s, left + Inches(0.35), top + Inches(0.16), cw - Inches(0.5), Inches(0.4),
                 [one(title, 14.5, WHITE, bold=True)])
        add_text(s, left + Inches(0.35), top + Inches(0.56), cw - Inches(0.5), Inches(0.5),
                 [one(desc, 11.5, GRAY)], line_spacing=1.05)

    # ---- Slide 5: Tech Stack ----------------------------------------------
    s = new_slide()
    header(s, "Under the Hood", "Technology Stack", 5)
    columns = [
        ("FRONTEND", BLUE, [
            "React 18", "Create React App", "Component-based UI",
            "Fetch API → REST", "localStorage state"]),
        ("BACKEND", GREEN, [
            "FastAPI + Uvicorn", "SQLAlchemy (engine)", "Pydantic models",
            "Async query pipeline", "Layered architecture"]),
        ("AI / ML", PURPLE, [
            "Azure OpenAI GPT-4.1", "text-embedding-3-small", "LangChain",
            "FAISS vector store", "RAG retrieval"]),
        ("DATA", AMBER, [
            "PostgreSQL", "psycopg2 driver", "pandas (file parse)",
            "Foreign-key graph", "CSV / Excel ingest"]),
    ]
    cw = Inches(2.85)
    gap = Inches(0.23)
    total = 4 * cw + 3 * gap
    start = (EMU_W - total) / 2
    top = Inches(2.05)
    height = Inches(4.35)
    for i, (title, col, rows) in enumerate(columns):
        left = start + i * (cw + gap)
        add_rect(s, left, top, cw, height, BG_CARD, rounded=True, radius=0.05)
        head = add_rect(s, left, top, cw, Inches(0.7), col, rounded=True, radius=0.12)
        add_text(s, left, top, cw, Inches(0.7), [one(title, 15, BG_DARK, bold=True)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y = top + Inches(0.95)
        for r in rows:
            add_rect(s, left + Inches(0.28), y + Inches(0.07), Inches(0.13), Inches(0.13), col,
                     rounded=True, radius=0.5)
            add_text(s, left + Inches(0.52), y, cw - Inches(0.7), Inches(0.4),
                     [one(r, 12.5, WHITE)])
            y += Inches(0.62)

    # ---- Slide 6: System Architecture -------------------------------------
    s = new_slide()
    header(s, "Design", "System Architecture", 6)

    def layer(left, top, width, height, title, subtitle, color, title_color=WHITE):
        box = add_rect(s, left, top, width, height, BG_CARD, rounded=True, radius=0.08)
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        add_rect(s, left, top, Inches(0.14), height, color)
        add_text(s, left + Inches(0.32), top + Inches(0.12), width - Inches(0.5), Inches(0.45),
                 [one(title, 15, title_color, bold=True)])
        add_text(s, left + Inches(0.32), top + Inches(0.52), width - Inches(0.5), Inches(0.5),
                 [one(subtitle, 11, GRAY)], line_spacing=1.0)
        return box

    # Frontend
    layer(Inches(2.3), Inches(1.85), Inches(8.7), Inches(0.95), "React Frontend (port 3000)",
          "Query box · Schema browser · Data viewer · File upload · Results panel", BLUE)
    # arrow down
    a1 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.35), Inches(2.86), Inches(0.6), Inches(0.42))
    fill_shape(a1, TEAL)
    add_text(s, Inches(7.0), Inches(2.9), Inches(4.0), Inches(0.35),
             [one("REST / JSON  (HTTP)", 10.5, GRAY, bold=True)])

    # Backend container
    bc = add_outline_rect(s, Inches(1.0), Inches(3.45), Inches(8.2), Inches(2.75), TEAL_DK, weight=1.5)
    add_text(s, Inches(1.2), Inches(3.5), Inches(6), Inches(0.4),
             [one("FastAPI Backend (port 8000)", 13, TEAL, bold=True)])
    sub = [
        ("Routes", "query · upload\ndatabases · metadata", GREEN),
        ("Query Pipeline", "orchestrates the\n5-step flow", TEAL),
        ("Core", "embedder · SQL gen\nexecutor · files", PURPLE),
        ("Utils", "join graph · prompt\nmetadata · state", AMBER),
    ]
    sw = Inches(1.85)
    sx = Inches(1.22)
    for i, (t, d, col) in enumerate(sub):
        left = sx + i * (sw + Inches(0.12))
        card = add_rect(s, left, Inches(4.0), sw, Inches(2.0), BG_CARD2, rounded=True, radius=0.08)
        add_rect(s, left, Inches(4.0), sw, Inches(0.1), col)
        add_text(s, left + Inches(0.12), Inches(4.2), sw - Inches(0.24), Inches(0.4),
                 [one(t, 13, col, bold=True)], align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.1), Inches(4.75), sw - Inches(0.2), Inches(1.1),
                 [one(d, 10.5, GRAY)], align=PP_ALIGN.CENTER, line_spacing=1.05)

    # Right side external services
    ext = [("Azure OpenAI", "GPT-4.1 + embeddings", PURPLE, Inches(3.5)),
           ("FAISS", "vector similarity", AMBER, Inches(4.55)),
           ("PostgreSQL", "data + schema", TEAL, Inches(5.6))]
    for name, desc, col, top in ext:
        card = add_rect(s, Inches(9.55), top, Inches(3.0), Inches(0.92), BG_CARD, rounded=True, radius=0.12)
        card.line.color.rgb = col
        card.line.width = Pt(1.25)
        add_text(s, Inches(9.75), top + Inches(0.13), Inches(2.7), Inches(0.4),
                 [one(name, 13, col, bold=True)])
        add_text(s, Inches(9.75), top + Inches(0.5), Inches(2.7), Inches(0.35),
                 [one(desc, 10.5, GRAY)])
    # connector backend -> services
    conn = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, Inches(9.2), Inches(4.6), Inches(0.4), Inches(0.34))
    fill_shape(conn, TEAL)

    # ---- Slide 7: Query Pipeline ------------------------------------------
    s = new_slide()
    header(s, "Core Flow", "The Query Pipeline", 7)
    add_text(s, Inches(0.85), Inches(1.7), Inches(11.5), Inches(0.5),
             [one("Every question travels through five orchestrated steps — from embedding to answer.", 13.5, GRAY)])
    steps = [
        ("1", "Embed & Retrieve", "Embed the question and run FAISS semantic search to find the top relevant tables.", BLUE),
        ("2", "Gather Context", "Collect metadata for those tables plus any uploaded tables.", PURPLE),
        ("3", "Load Join Graph", "Pull the database's foreign-key join graph for correct joins.", TEAL),
        ("4", "Generate SQL", "GPT-4.1 writes the SQL and flags dangerous operations in one call.", AMBER),
        ("5", "Validate & Execute", "Validate, run read-only against PostgreSQL, then AI-summarize the rows.", GREEN),
    ]
    n = len(steps)
    cw = Inches(2.28)
    gap = Inches(0.12)
    total = n * cw + (n - 1) * gap
    start = (EMU_W - total) / 2
    top = Inches(2.55)
    h = Inches(3.3)
    for i, (num, title, desc, col) in enumerate(steps):
        left = start + i * (cw + gap)
        add_rect(s, left, top, cw, h, BG_CARD, rounded=True, radius=0.07)
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + cw/2 - Inches(0.42), top + Inches(0.3),
                                  Inches(0.84), Inches(0.84))
        fill_shape(circ, col)
        add_text(s, left + cw/2 - Inches(0.42), top + Inches(0.3), Inches(0.84), Inches(0.84),
                 [one(num, 26, BG_DARK, bold=True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.12), top + Inches(1.35), cw - Inches(0.24), Inches(0.7),
                 [one(title, 14, WHITE, bold=True)], align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.18), top + Inches(2.0), cw - Inches(0.36), Inches(1.2),
                 [one(desc, 11, GRAY)], align=PP_ALIGN.CENTER, line_spacing=1.1)
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, left + cw - Inches(0.02), top + h/2 - Inches(0.18),
                                    Inches(0.22), Inches(0.36))
            fill_shape(ar, TEAL)

    # ---- Slide 8: Semantic Retrieval (RAG) --------------------------------
    s = new_slide()
    header(s, "AI Deep-Dive", "Semantic Table Retrieval (RAG)", 8)
    add_text(s, Inches(0.85), Inches(1.75), Inches(11.4), Inches(0.6),
             [one("Instead of dumping the whole schema into the prompt, we retrieve only the tables that "
                  "matter — faster, cheaper and more accurate.", 13.5, GRAY)], line_spacing=1.2)

    flow = [
        ("Question", "Free-text query", BLUE),
        ("Embedding", "text-embedding-3-small", PURPLE),
        ("FAISS Search", "IndexFlatL2 · top-k=5", AMBER),
        ("Relevant Tables", "L2 distance ≤ 2.5", GREEN),
    ]
    bw = Inches(2.6)
    gap = Inches(0.55)
    total = 4 * bw + 3 * gap
    start = (EMU_W - total) / 2
    top = Inches(2.75)
    for i, (t, d, col) in enumerate(flow):
        left = start + i * (bw + gap)
        card = add_rect(s, left, top, bw, Inches(1.5), BG_CARD, rounded=True, radius=0.1)
        card.line.color.rgb = col
        card.line.width = Pt(1.25)
        add_text(s, left, top + Inches(0.3), bw, Inches(0.5), [one(t, 16, WHITE, bold=True)],
                 align=PP_ALIGN.CENTER)
        add_text(s, left, top + Inches(0.85), bw, Inches(0.5), [one(d, 11, col, bold=True)],
                 align=PP_ALIGN.CENTER)
        if i < 3:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + bw + Inches(0.07),
                                    top + Inches(0.6), Inches(0.4), Inches(0.3))
            fill_shape(ar, TEAL)

    # bottom note cards
    notes = [
        ("Startup embedding", "All tables embedded concurrently on boot (bounded by a semaphore).", TEAL),
        ("On disk", "Indices saved as *_index.faiss with *_metadata.json companions.", AMBER),
        ("Async & parallel", "Async Azure client embeds many tables at once for speed.", PURPLE),
    ]
    nx = Inches(0.85)
    nw = Inches(3.78)
    for i, (t, d, col) in enumerate(notes):
        left = nx + i * (nw + Inches(0.2))
        add_rect(s, left, Inches(4.85), nw, Inches(1.5), BG_CARD, rounded=True, radius=0.08)
        add_rect(s, left, Inches(4.85), nw, Inches(0.1), col)
        add_text(s, left + Inches(0.25), Inches(5.05), nw - Inches(0.5), Inches(0.4),
                 [one(t, 13.5, col, bold=True)])
        add_text(s, left + Inches(0.25), Inches(5.5), nw - Inches(0.5), Inches(0.8),
                 [one(d, 11, GRAY)], line_spacing=1.15)

    # ---- Slide 9: SQL Generation & Safety ---------------------------------
    s = new_slide()
    header(s, "AI Deep-Dive", "SQL Generation & Safety", 9)
    # left: prompt construction
    add_text(s, Inches(0.85), Inches(1.8), Inches(6.0), Inches(0.4),
             [one("What goes into the prompt", 16, TEAL, bold=True)])
    promptitems = [
        ("Question", "The user's natural-language request.", BLUE),
        ("Table schemas", "Columns, types, nullability, primary keys.", PURPLE),
        ("Join paths", "Foreign-key relationships between tables.", TEAL),
        ("Rules", "SELECT-only, pick specific columns, 2-dp numerics.", AMBER),
    ]
    bullet_block(s, Inches(0.85), Inches(2.45), Inches(6.0), promptitems, gap=0.85, title_size=14, desc_size=11)

    # right: safety panel
    rcard = add_rect(s, Inches(7.2), Inches(1.85), Inches(5.3), Inches(4.5), BG_CARD, rounded=True, radius=0.06)
    add_rect(s, Inches(7.2), Inches(1.85), Inches(5.3), Inches(0.62), RED, rounded=True, radius=0.2)
    add_text(s, Inches(7.2), Inches(1.85), Inches(5.3), Inches(0.62),
             [one("⚠  SAFETY GUARDRAILS", 14, WHITE, bold=True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    safety = [
        ("Single-call danger check", "LLM returns DANGEROUS / REASON / SQL together.", RED),
        ("Blocked, not run", "DELETE, DROP, UPDATE, INSERT are shown but never executed.", RED),
        ("Query validation", "Syntax checked before it reaches the database.", AMBER),
        ("Read-only + LIMIT", "Auto-applies a row LIMIT; results capped for safety.", GREEN),
    ]
    bullet_block(s, Inches(7.45), Inches(2.75), Inches(4.85), safety, gap=0.85, title_size=13, desc_size=10.5)

    # ---- Slide 10: Join Graph ---------------------------------------------
    s = new_slide()
    header(s, "AI Deep-Dive", "Join-Graph Intelligence", 10)
    add_text(s, Inches(0.85), Inches(1.75), Inches(11.3), Inches(0.7),
             [one("At startup the backend inspects every table's foreign keys and builds an adjacency "
                  "list of joins per database — so multi-table questions are answered correctly.", 13.5, GRAY)],
             line_spacing=1.2)

    # mini ER-style diagram
    def ent(left, top, name, col):
        box = add_rect(s, left, top, Inches(2.2), Inches(1.0), BG_CARD, rounded=True, radius=0.1)
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        add_rect(s, left, top, Inches(2.2), Inches(0.34), col, rounded=True, radius=0.2)
        add_text(s, left, top, Inches(2.2), Inches(0.34), [one(name, 12.5, BG_DARK, bold=True)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return box

    ent(Inches(1.2), Inches(3.0), "customers", BLUE)
    ent(Inches(5.5), Inches(2.55), "orders", TEAL)
    ent(Inches(5.5), Inches(3.9), "products", PURPLE)
    ent(Inches(9.8), Inches(3.0), "order_items", AMBER)
    # connectors
    for (x1, y1, w, h) in [(3.4, 3.5, 2.1, 0.05), (7.7, 3.05, 2.1, 0.05)]:
        ln = add_rect(s, Inches(x1), Inches(y1), Inches(w), Inches(0.04), TEAL)
    add_text(s, Inches(3.4), Inches(3.15), Inches(2.1), Inches(0.3),
             [one("customer_id", 9.5, GRAY, bold=True)], align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.7), Inches(2.7), Inches(2.1), Inches(0.3),
             [one("order_id", 9.5, GRAY, bold=True)], align=PP_ALIGN.CENTER)

    # bottom explanation cards
    jb = [
        ("Built once at boot", "SQLAlchemy inspector reads FK constraints for all databases.", TEAL),
        ("Filtered per query", "Only joins linking the retrieved tables enter the prompt.", PURPLE),
        ("Path extraction", "The actual join path used is parsed back from the final SQL.", AMBER),
    ]
    nx = Inches(0.85)
    nw = Inches(3.78)
    for i, (t, d, col) in enumerate(jb):
        left = nx + i * (nw + Inches(0.2))
        add_rect(s, left, Inches(4.9), nw, Inches(1.5), BG_CARD, rounded=True, radius=0.08)
        add_rect(s, left, Inches(4.9), nw, Inches(0.1), col)
        add_text(s, left + Inches(0.25), Inches(5.1), nw - Inches(0.5), Inches(0.4),
                 [one(t, 13, col, bold=True)])
        add_text(s, left + Inches(0.25), Inches(5.52), nw - Inches(0.5), Inches(0.8),
                 [one(d, 11, GRAY)], line_spacing=1.15)

    # ---- Slide 11: File Upload --------------------------------------------
    s = new_slide()
    header(s, "Feature", "Upload & Query Your Own Data", 11)
    add_text(s, Inches(0.85), Inches(1.8), Inches(11.2), Inches(0.6),
             [one("Bring a CSV or Excel file and query it immediately alongside the database tables.", 13.5, GRAY)])
    upflow = [
        ("Upload", "CSV / Excel via the UI", BLUE),
        ("Parse", "pandas reads the file", PURPLE),
        ("Create table", "Temp table in PostgreSQL", AMBER),
        ("Query", "Always in LLM context", GREEN),
    ]
    bw = Inches(2.6)
    gap = Inches(0.55)
    total = 4 * bw + 3 * gap
    start = (EMU_W - total) / 2
    top = Inches(2.7)
    for i, (t, d, col) in enumerate(upflow):
        left = start + i * (bw + gap)
        card = add_rect(s, left, top, bw, Inches(1.55), BG_CARD, rounded=True, radius=0.1)
        card.line.color.rgb = col
        card.line.width = Pt(1.25)
        add_text(s, left, top + Inches(0.32), bw, Inches(0.5), [one(t, 16, WHITE, bold=True)],
                 align=PP_ALIGN.CENTER)
        add_text(s, left, top + Inches(0.9), bw, Inches(0.5), [one(d, 11, col, bold=True)],
                 align=PP_ALIGN.CENTER)
        if i < 3:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + bw + Inches(0.07),
                                    top + Inches(0.62), Inches(0.4), Inches(0.3))
            fill_shape(ar, TEAL)
    upnotes = [
        ("Clean & safe naming", "Columns normalized; tables prefixed uploaded__ with a UUID.", TEAL),
        ("Auto-cleanup", "Leftover uploaded tables are removed on each restart.", AMBER),
        ("Shown under its DB", "Uploaded files appear with a 📤 icon under their database.", PURPLE),
    ]
    nx = Inches(0.85)
    nw = Inches(3.78)
    for i, (t, d, col) in enumerate(upnotes):
        left = nx + i * (nw + Inches(0.2))
        add_rect(s, left, Inches(4.85), nw, Inches(1.5), BG_CARD, rounded=True, radius=0.08)
        add_rect(s, left, Inches(4.85), nw, Inches(0.1), col)
        add_text(s, left + Inches(0.25), Inches(5.05), nw - Inches(0.5), Inches(0.4),
                 [one(t, 13, col, bold=True)])
        add_text(s, left + Inches(0.25), Inches(5.48), nw - Inches(0.5), Inches(0.8),
                 [one(d, 11, GRAY)], line_spacing=1.15)

    # ---- Slide 12: Backend Architecture -----------------------------------
    s = new_slide()
    header(s, "Code Structure", "Backend Architecture", 12)
    add_text(s, Inches(0.85), Inches(1.75), Inches(11.3), Inches(0.5),
             [one("A clean, layered FastAPI codebase — thin routes delegate to a service pipeline.", 13.5, GRAY)])
    modules = [
        ("routes/", "HTTP endpoints — query, upload, databases, metadata, join_graphs", GREEN),
        ("services/", "query_pipeline.py — orchestrates the end-to-end flow", TEAL),
        ("core/", "embedder · vector_store · sql_generator · query_executor · file_handler", PURPLE),
        ("config/", "azure_client.py & llm.py — Azure OpenAI + LangChain setup", BLUE),
        ("db/", "db.py engine factory · db_models.py Pydantic request/response models", AMBER),
        ("utils/", "join_graph_builder · metadata_builder · prompt_builder · state · logger", PINK),
    ]
    gx, gy = Inches(0.85), Inches(2.4)
    cw, ch = Inches(5.75), Inches(1.18)
    for i, (name, desc, col) in enumerate(modules):
        col_i = i % 2
        row_i = i // 2
        left = gx + col_i * (cw + Inches(0.3))
        top = gy + row_i * (ch + Inches(0.12))
        add_rect(s, left, top, cw, ch, BG_CARD, rounded=True, radius=0.09)
        add_rect(s, left, top, Inches(0.12), ch, col)
        add_text(s, left + Inches(0.35), top + Inches(0.18), cw - Inches(0.5), Inches(0.4),
                 [one(name, 15, col, bold=True, font=MONO)])
        add_text(s, left + Inches(0.35), top + Inches(0.6), cw - Inches(0.55), Inches(0.5),
                 [one(desc, 11, GRAY)], line_spacing=1.05)

    # ---- Slide 13: Frontend Architecture ----------------------------------
    s = new_slide()
    header(s, "Code Structure", "Frontend Architecture", 13)
    add_text(s, Inches(0.85), Inches(1.75), Inches(11.3), Inches(0.5),
             [one("A React 18 single-page app with feature-grouped components and three main views.", 13.5, GRAY)])
    # three views
    views = [
        ("Query View", "Ask questions, see SQL, results, summary, complexity & token usage.", BLUE),
        ("Schema View", "Browse databases, tables and columns; uploaded files shown inline.", PURPLE),
        ("Data View", "Inspect raw table data with a full-screen table modal.", TEAL),
    ]
    vw = Inches(3.78)
    for i, (t, d, col) in enumerate(views):
        left = Inches(0.85) + i * (vw + Inches(0.2))
        add_rect(s, left, Inches(2.35), vw, Inches(1.7), BG_CARD, rounded=True, radius=0.08)
        add_rect(s, left, Inches(2.35), vw, Inches(0.5), col, rounded=True, radius=0.18)
        add_text(s, left, Inches(2.35), vw, Inches(0.5), [one(t, 14, BG_DARK, bold=True)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.22), Inches(3.0), vw - Inches(0.44), Inches(0.9),
                 [one(d, 11.5, GRAY)], line_spacing=1.2)

    add_text(s, Inches(0.85), Inches(4.35), Inches(11), Inches(0.4),
             [one("Component groups", 14, TEAL, bold=True)])
    groups = ["layout/ — Sidebar", "query/ — Results, Modal, ChatMessage", "schema/ — DatabaseSchemaView",
              "data/ — ViewDataTable", "upload/ — FileUploadModal", "common/ — FullScreenTableModal"]
    gx = Inches(0.85)
    gy = Inches(4.95)
    gw = Inches(3.78)
    for i, g in enumerate(groups):
        col_i = i % 3
        row_i = i // 3
        left = gx + col_i * (gw + Inches(0.2))
        top = gy + row_i * Inches(0.72)
        card = add_rect(s, left, top, gw, Inches(0.58), BG_CARD2, rounded=True, radius=0.18)
        add_rect(s, left + Inches(0.18), top + Inches(0.21), Inches(0.14), Inches(0.14), TEAL,
                 rounded=True, radius=0.5)
        add_text(s, left + Inches(0.45), top, gw - Inches(0.6), Inches(0.58),
                 [one(g, 11.5, WHITE)], anchor=MSO_ANCHOR.MIDDLE)

    # ---- Slide 14: API Endpoints ------------------------------------------
    s = new_slide()
    header(s, "Interface", "REST API Endpoints", 14)
    endpoints = [
        ("POST", "/query", "Natural-language query → SQL, results & summary", GREEN),
        ("GET", "/query/search-tables", "Semantic table search via embeddings", BLUE),
        ("GET", "/databases", "List all accessible PostgreSQL databases", TEAL),
        ("GET", "/databases/{db}/tables", "List tables (with row counts) in a database", TEAL),
        ("GET", "/databases/{db}/schema", "Full schema for every table in a database", PURPLE),
        ("POST", "/upload", "Upload a CSV / Excel file as a queryable table", AMBER),
        ("GET", "/upload/tables", "List uploaded tables for a database", AMBER),
        ("GET", "/metadata  ·  /join-graphs", "Cached metadata & foreign-key join graphs", PINK),
    ]
    top = Inches(1.8)
    rowh = Inches(0.56)
    for i, (verb, path, desc, col) in enumerate(endpoints):
        y = top + i * (rowh + Inches(0.035))
        add_rect(s, Inches(0.85), y, Inches(11.6), rowh, BG_CARD, rounded=True, radius=0.15)
        # verb badge
        vb = add_rect(s, Inches(1.0), y + Inches(0.1), Inches(0.95), Inches(0.36), col, rounded=True, radius=0.25)
        add_text(s, Inches(1.0), y + Inches(0.1), Inches(0.95), Inches(0.36),
                 [one(verb, 11, BG_DARK, bold=True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.15), y, Inches(4.2), rowh, [one(path, 13, WHITE, bold=True, font=MONO)],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.5), y, Inches(5.8), rowh, [one(desc, 12, GRAY)], anchor=MSO_ANCHOR.MIDDLE)

    # ---- Slide 15: Security -----------------------------------------------
    s = new_slide()
    header(s, "Trust", "Security & Safety", 15)
    sec = [
        ("Destructive-op blocking", "DELETE / DROP / UPDATE / INSERT are detected and never executed.", RED),
        ("Read-only execution", "Queries run with an enforced row LIMIT for safe previews.", GREEN),
        ("Query validation", "Basic syntax validation before hitting the database.", AMBER),
        ("Secrets in env", "Azure keys & DB credentials live in .env.local, never committed.", BLUE),
        ("Parameterized engine", "SQLAlchemy engine per database with managed connections.", PURPLE),
        ("Session cleanup", "Uploaded tables & embeddings are cleaned up on shutdown/restart.", TEAL),
    ]
    gx, gy = Inches(0.85), Inches(2.0)
    cw, ch = Inches(5.75), Inches(1.35)
    for i, (title, desc, col) in enumerate(sec):
        col_i = i % 2
        row_i = i // 2
        left = gx + col_i * (cw + Inches(0.3))
        top = gy + row_i * (ch + Inches(0.15))
        add_rect(s, left, top, cw, ch, BG_CARD, rounded=True, radius=0.08)
        # shield dot
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), top + Inches(0.3), Inches(0.5), Inches(0.5))
        fill_shape(circ, col)
        add_text(s, left + Inches(0.25), top + Inches(0.3), Inches(0.5), Inches(0.5),
                 [one("✓", 18, BG_DARK, bold=True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.95), top + Inches(0.2), cw - Inches(1.1), Inches(0.45),
                 [one(title, 14.5, WHITE, bold=True)])
        add_text(s, left + Inches(0.95), top + Inches(0.62), cw - Inches(1.15), Inches(0.6),
                 [one(desc, 11, GRAY)], line_spacing=1.1)

    # ---- Slide 16: Roadmap ------------------------------------------------
    s = new_slide()
    header(s, "What's Next", "Future Roadmap", 16)
    road = [
        ("Query history persistence", "Save and revisit past questions and results.", BLUE),
        ("User authentication", "Per-user access and database permissions.", PURPLE),
        ("Data export", "Download results as CSV / JSON.", GREEN),
        ("Optimization hints", "Suggest indexes and query improvements.", AMBER),
        ("Charts & visualizations", "Auto-render graphs from result sets.", TEAL),
        ("Multi-turn conversations", "Follow-up questions with shared context.", PINK),
    ]
    gx, gy = Inches(0.85), Inches(2.05)
    cw, ch = Inches(5.75), Inches(1.3)
    for i, (title, desc, col) in enumerate(road):
        col_i = i % 2
        row_i = i // 2
        left = gx + col_i * (cw + Inches(0.3))
        top = gy + row_i * (ch + Inches(0.18))
        add_rect(s, left, top, cw, ch, BG_CARD, rounded=True, radius=0.08)
        add_rect(s, left, top, Inches(0.12), ch, col)
        # arrow icon
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(0.32), top + Inches(0.46),
                                Inches(0.42), Inches(0.34))
        fill_shape(ar, col)
        add_text(s, left + Inches(0.95), top + Inches(0.2), cw - Inches(1.1), Inches(0.45),
                 [one(title, 14.5, WHITE, bold=True)])
        add_text(s, left + Inches(0.95), top + Inches(0.62), cw - Inches(1.15), Inches(0.55),
                 [one(desc, 11, GRAY)], line_spacing=1.1)

    # ---- Slide 17: Thank You ----------------------------------------------
    s = new_slide()
    add_rect(s, Inches(0), Inches(0), Inches(0.28), EMU_H, TEAL)
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(3.8), Inches(5.0), Inches(5.0))
    fill_shape(c1, BG_CARD)
    c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.6), Inches(4.6), Inches(4.6))
    fill_shape(c2, BG_CARD2)
    hero_logo(s)

    add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.3),
             [[("Thank ", 60, WHITE, True, FONT), ("You", 60, TEAL, True, FONT)]])
    add_text(s, Inches(0.95), Inches(3.9), Inches(10.8), Inches(0.7),
             [one("Text-to-SQL — making databases speak your language.", 20, GRAY)])
    x = Inches(0.95)
    for label, col in [("React", BLUE), ("FastAPI", GREEN), ("Azure OpenAI", PURPLE),
                       ("FAISS", AMBER), ("PostgreSQL", TEAL)]:
        w = chip(s, x, Inches(5.0), label, col)
        x += w + Inches(0.18)
    add_text(s, Inches(0.95), Inches(6.4), Inches(11), Inches(0.5),
             [one("Questions & Discussion", 14, GRAY, bold=True)])
    add_footer(s)

    # Save
    out = Path(__file__).resolve().parent / "TextToSQL_Presentation.pptx"
    prs.save(str(out))
    print(f"Saved: {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
